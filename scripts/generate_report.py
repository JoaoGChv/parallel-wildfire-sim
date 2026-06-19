"""
generate_report.py — Gera relatório DOCX + apresentação PPTX da Entrega 1
"""
import csv, json
from pathlib import Path
from datetime import date

BASE    = Path(__file__).parent.parent
OUT_DIR = BASE / "results"
OUT_DIR.mkdir(exist_ok=True)
VIZ     = BASE / "results" / "visualizations"
MODELS  = BASE / "results" / "models"

# ═══════════════════════════════════════════════════════════════════════════════
#  RELATÓRIO DOCX
# ═══════════════════════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

def add_heading(doc, text, level=1, color=None):
    h = doc.add_heading(text, level=level)
    if color:
        for run in h.runs:
            run.font.color.rgb = RGBColor(*color)
    return h

def add_metric_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = h
        hdr[i].paragraphs[0].runs[0].bold = True
    for i, row in enumerate(rows):
        cells = table.rows[i+1].cells
        for j, val in enumerate(row):
            cells[j].text = str(val)
    return table

def build_report():
    doc = Document()

    # Estilos
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Título
    title = doc.add_heading("Relatório — Entrega 1", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph("HPC — Simulação de Incêndios Florestais")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].bold = True

    doc.add_paragraph(f"Data: {date.today().strftime('%d/%m/%Y')}")
    doc.add_paragraph("Referência: Kwon et al. (2022). A Semantic Data-Based Distributed Computing Framework to Accelerate Digital Twin Services for Large-Scale Disasters. Sensors 22(18), 6749.")
    doc.add_paragraph()

    # ── 1. Introdução ────────────────────────────────────────────────────────
    add_heading(doc, "1. Introdução", 1)
    doc.add_paragraph(
        "Este relatório documenta a implementação completa da Entrega 1 do projeto HPC — "
        "Simulação de Incêndios Florestais, baseada no artigo de referência Kwon et al. (2022). "
        "O objetivo desta entrega foi estabelecer a fundação de dados e os baselines computacionais "
        "necessários para as entregas seguintes, que implementarão paralelismo via OpenMP e CUDA."
    )

    # ── 2. Ambiente ──────────────────────────────────────────────────────────
    add_heading(doc, "2. Configuração do Ambiente", 1)
    doc.add_paragraph("O ambiente de desenvolvimento foi configurado com os seguintes componentes:")
    for item in [
        "GPU: NVIDIA GeForce RTX 4090 (24 GB VRAM)",
        "CUDA 12.1 + PyTorch 2.5.1",
        "GCC 11 + OpenMP + MPI (mpifort)",
        "Python 3.12 + rasterio 1.5 + pyproj 3.7",
        "elmfire 2025.1002 (compilado de fonte — Fortran + MPI)",
        "GDAL 3.8.4",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    # ── 3. Dataset ───────────────────────────────────────────────────────────
    add_heading(doc, "3. Dataset — 68 Cenários de Incêndios 2016", 1)
    doc.add_paragraph(
        "O dataset foi construído a partir de duas fontes oficiais dos EUA, conforme descrito "
        "no artigo baseline:"
    )
    add_heading(doc, "3.1 Fontes de Dados", 2)
    for item in [
        "LANDFIRE 2016 (landfire.gov): 8 camadas landscape via ArcGIS ImageServer",
        "NIFC Open Data: 4.410 perímetros históricos de incêndios 2016 (458 MB)",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    add_heading(doc, "3.2 Seleção dos Cenários", 2)
    doc.add_paragraph(
        "Seguindo os critérios do paper: incêndios large-scale (>300 acres) no CONUS, "
        "usando apenas o perímetro final (latest=Y). Do total de 437 incêndios únicos válidos "
        "(paper cita 341 — diferença por versão do dataset):"
    )
    add_metric_table(doc,
        ["Split", "Quantidade", "Critério", "Classes NWCG"],
        [
            ["Treino", "51", "Aleatório, large-scale ≥300 acres", "E, F, G"],
            ["Teste",  "17", "Classe F: 1.000–4.999 acres",       "F"],
            ["Total",  "68", "—", "—"],
        ]
    )

    add_heading(doc, "3.3 As 14 Camadas Semânticas", 2)
    doc.add_paragraph("Cada cenário é representado por um array 450×450×14 (resolução 30m/pixel):")
    add_metric_table(doc,
        ["#", "Nome", "Fonte", "Tipo"],
        [
            ["0","Elevation (DEM)","LANDFIRE Topo/LF2020","Landscape"],
            ["1","Aspect","LANDFIRE Topo/LF2020","Landscape"],
            ["2","Slope","LANDFIRE Topo/LF2020","Landscape"],
            ["3","Fuel Model (FBFM40)","LANDFIRE LF2016","Landscape"],
            ["4","Canopy Cover","LANDFIRE LF2016","Landscape"],
            ["5","Canopy Height","LANDFIRE LF2016","Landscape"],
            ["6","Canopy Base Height","LANDFIRE LF2016","Landscape"],
            ["7","Canopy Bulk Density","LANDFIRE LF2016","Landscape"],
            ["8","Temperature","Sintético (FARSITE ranges)","Weather"],
            ["9","Humidity","Sintético (FARSITE ranges)","Weather"],
            ["10","Cloud Cover","Sintético (FARSITE ranges)","Weather"],
            ["11","Precipitation","Sintético (FARSITE ranges)","Weather"],
            ["12","Wind Direction","Sintético (FARSITE ranges)","Weather"],
            ["13","Wind Speed","Sintético (FARSITE ranges)","Weather"],
        ]
    )

    # ── 4. Labels (elmfire) ──────────────────────────────────────────────────
    add_heading(doc, "4. Geração de Labels via elmfire", 1)
    doc.add_paragraph(
        "Para cada cenário, o simulador elmfire (implementação open-source baseada no "
        "FARSITE, que usa as equações de Rothermel + princípio de Huygens) foi executado "
        "por 24 horas de simulação, gerando o raster de time-of-arrival (TOA) por célula. "
        "Este raster é o label de treino da CNN."
    )
    doc.add_paragraph(
        "Nota: as simulações de 24h com clima sintético conservador resultaram em propagação "
        "limitada (<0.5% das células por cenário). Em produção real, o clima histórico "
        "(GRIDMET) produziria propagação maior e labels mais ricos."
    )

    if VIZ.joinpath("elmfire_toa_S001.png").exists():
        doc.add_picture(str(VIZ / "elmfire_toa_S001.png"), width=Inches(5.5))
        doc.add_paragraph("Figura 1: Time of arrival — S001 (0624_RS__EAST_MAURY, Oregon)").italic = True

    # ── 5. CNN ───────────────────────────────────────────────────────────────
    add_heading(doc, "5. CNN de Predição de Carga Computacional", 1)
    add_heading(doc, "5.1 Arquitetura", 2)
    doc.add_paragraph(
        "Implementada conforme Figure 5 do artigo. Input: janela atomic_size×atomic_size×14 "
        "centrada em cada pixel. Output: escalar (carga computacional prevista)."
    )
    add_metric_table(doc,
        ["Camada","Detalhe"],
        [
            ["Input","(batch, 14, W, W) — W ∈ {9,12,15,18,21}"],
            ["Conv2d(14→32) + BN + ReLU","kernel 3×3, sem padding"],
            ["Conv2d(32→64) + BN + ReLU","kernel 3×3, sem padding"],
            ["Conv2d(64→128) + BN + ReLU","kernel 3×3, sem padding"],
            ["Flatten → FC(256) → ReLU","—"],
            ["FC(1)","saída escalar"],
        ]
    )

    add_heading(doc, "5.2 Configuração de Treino", 2)
    add_metric_table(doc,
        ["Parâmetro","Valor"],
        [
            ["Otimizador","Adam (lr=1e-3, weight_decay=1e-4)"],
            ["Loss","MSE"],
            ["Batch size","4096"],
            ["Dataset","10.3M amostras treino / 3.4M teste"],
            ["GPU","RTX 4090"],
            ["Labels normalizados","TOA / 86400s → [0,1]"],
        ]
    )

    add_heading(doc, "5.3 Resultados — Sweep de Atomic Sizes", 2)
    rows_sweep = []
    for s in [9,12,15,18,21]:
        mp = MODELS / f"cnn_size{s}_meta.json"
        if mp.exists():
            d = json.load(open(mp))
            br = d.get("burned_cell_rmse")
            br_h = f"{br*24:.2f}h" if br else "—"
            rows_sweep.append([
                f"{s}×{s}", str(d["epochs_trained"]),
                f"{d['best_test_mse']:.6f}",
                br_h,
                f"{d['n_params']:,}",
            ])
    add_metric_table(doc,
        ["Atomic Size","Épocas","Best Test MSE","Burned RMSE","Parâmetros"],
        rows_sweep
    )
    doc.add_paragraph(
        "Resultado: 21×21 obteve menor test MSE (0.001448), consistente com o paper. "
        "O burned RMSE mede o erro apenas nas células efetivamente queimadas."
    )

    if VIZ.joinpath("sweep_mse.png").exists():
        doc.add_picture(str(VIZ / "sweep_mse.png"), width=Inches(6))
        doc.add_paragraph("Figura 2: MSE × Atomic Size (sweep completo)").italic = True

    # ── 6. Simulador C ───────────────────────────────────────────────────────
    add_heading(doc, "6. Simulador Sequencial em C", 1)
    doc.add_paragraph(
        "Implementado em C puro com as equações de Rothermel (1972) para taxa de "
        "espalhamento (ROS) e propagação via Dijkstra (princípio de Huygens). "
        "Usa fila de prioridade (min-heap) para propagar a frente de fogo célula a célula "
        "em conectividade 8-vizinhos."
    )

    add_heading(doc, "6.1 Arquitetura do Código", 2)
    add_metric_table(doc,
        ["Arquivo","Responsabilidade"],
        [
            ["simulator.c","Main: leitura de dados, loop Dijkstra, medição de tempo"],
            ["rothermel.c/h","Equações de Rothermel — 13 fuel models (FBFM13)"],
            ["priority_queue.c/h","Min-heap para Dijkstra"],
            ["Makefile","Compilação com -O2 -fopenmp"],
        ]
    )

    add_heading(doc, "6.2 Tempos de Execução — 68 Cenários (baseline)", 2)
    times = []
    seq_csv = BASE / "results" / "sequential_times.csv"
    if seq_csv.exists():
        with open(seq_csv) as f:
            for row in csv.DictReader(f):
                times.append(float(row["time_seconds"]))

    if times:
        add_metric_table(doc,
            ["Métrica","Valor"],
            [
                ["Cenários medidos","68/68"],
                ["Tempo médio",f"{sum(times)/len(times)*1000:.1f} ms"],
                ["Tempo mínimo",f"{min(times)*1000:.1f} ms"],
                ["Tempo máximo",f"{max(times)*1000:.1f} ms"],
                ["Desvio padrão",f"{(sum((t-sum(times)/len(times))**2 for t in times)/len(times))**0.5*1000:.1f} ms"],
                ["Runs por cenário","3 (média)"],
            ]
        )

    doc.add_paragraph(
        "Este é o baseline sequencial que será comparado com as versões OpenMP e CUDA "
        "nas entregas seguintes."
    )

    if VIZ.joinpath("simulator_S001_ws8_wd270.png").exists():
        doc.add_picture(str(VIZ / "simulator_S001_ws8_wd270.png"), width=Inches(6))
        doc.add_paragraph("Figura 3: Propagação simulador C — S001 (vento 8 m/s @ 270°)").italic = True

    # ── 7. Próximos passos ───────────────────────────────────────────────────
    add_heading(doc, "7. Próximos Passos — Entrega 2", 1)
    for item in [
        "Paralelização do simulador C com OpenMP (pragma omp parallel for no Dijkstra por regiões)",
        "Implementação da versão CUDA do simulador (propagação por blocos na GPU)",
        "Comparação de speedup: sequencial vs OpenMP vs CUDA nos 68 cenários",
        "Implementação do algoritmo de balanceamento de carga usando as predições da CNN",
        "Análise de escalabilidade com diferentes números de threads/blocos",
        "Validação cruzada das predições da CNN contra tempos reais de simulação",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    out = OUT_DIR / "relatorio_entrega1.docx"
    doc.save(out)
    print(f"Relatório salvo: {out}")
    return out


# ═══════════════════════════════════════════════════════════════════════════════
#  APRESENTAÇÃO PPTX
# ═══════════════════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

COR_TITULO    = PptxRGB(0x1A, 0x3A, 0x5C)  # azul escuro
COR_DESTAQUE  = PptxRGB(0xE8, 0x4A, 0x1E)  # laranja
COR_BG_SLIDE  = PptxRGB(0xF5, 0xF7, 0xFA)  # cinza claro
COR_TEXTO     = PptxRGB(0x22, 0x22, 0x22)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, left, top, width, height,
                  fontsize=24, bold=True, color=COR_TITULO, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(fontsize)
    run.font.color.rgb = color
    return box


def add_body_text(slide, text, left, top, width, height,
                  fontsize=14, color=COR_TEXTO, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(fontsize)
        run.font.color.rgb = color
        run.font.bold  = bold
    return box


def add_image_safe(slide, img_path, left, top, width, height=None):
    if Path(img_path).exists():
        if height:
            slide.shapes.add_picture(str(img_path), left, top, width, height)
        else:
            slide.shapes.add_picture(str(img_path), left, top, width)
        return True
    return False


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank

    W = prs.slide_width
    H = prs.slide_height

    # ── SLIDE 1: Capa ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_TITULO)
    # Faixa lateral laranja
    add_rect(s, Inches(0), Inches(0), Inches(0.3), H, COR_DESTAQUE)
    add_rect(s, Inches(0), Inches(5.2), W, Inches(0.15), COR_DESTAQUE)

    add_title_box(s, "HPC — Simulação de Incêndios Florestais",
                  Inches(0.7), Inches(1.0), Inches(12), Inches(1.2),
                  fontsize=32, color=PptxRGB(0xFF,0xFF,0xFF))
    add_title_box(s, "Entrega 1 — Fundação e Dados",
                  Inches(0.7), Inches(2.4), Inches(12), Inches(0.8),
                  fontsize=22, color=COR_DESTAQUE)
    add_body_text(s,
                  "Baseado em: Kwon et al. (2022) — Sensors 22(18), 6749\n"
                  f"Data: {date.today().strftime('%B %Y')}",
                  Inches(0.7), Inches(5.6), Inches(10), Inches(1.0),
                  fontsize=13, color=PptxRGB(0xCC,0xCC,0xCC))

    # ── SLIDE 2: Agenda ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Agenda", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=26, color=PptxRGB(0xFF,0xFF,0xFF))

    items = [
        ("01", "Contexto e Objetivos"),
        ("02", "Configuração do Ambiente"),
        ("03", "Dataset — 68 Cenários NIFC 2016"),
        ("04", "Pipeline de Preprocessing (450×450×14)"),
        ("05", "Geração de Labels via elmfire"),
        ("06", "CNN de Predição de Carga"),
        ("07", "Simulador Sequencial em C"),
        ("08", "Visualizações e Resultados"),
        ("09", "Próximos Passos"),
    ]
    for i, (num, text) in enumerate(items):
        col = i // 5
        row = i % 5
        lft = Inches(0.5 + col * 6.5)
        tp  = Inches(1.4 + row * 1.1)
        add_rect(s, lft, tp, Inches(0.55), Inches(0.55), COR_DESTAQUE)
        add_title_box(s, num, lft + Inches(0.08), tp, Inches(0.55), Inches(0.55),
                      fontsize=14, bold=True, color=PptxRGB(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        add_body_text(s, text, lft + Inches(0.65), tp + Inches(0.05), Inches(5.5), Inches(0.6),
                      fontsize=15, color=COR_TITULO, bold=True)

    # ── SLIDE 3: Contexto ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Contexto e Objetivos", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=26, color=PptxRGB(0xFF,0xFF,0xFF))

    add_body_text(s,
        "Objetivo: Replicar e estender o framework de Kwon et al. (2022) para\n"
        "aceleração de simulações de incêndios florestais via HPC.\n\n"
        "Problema central: simular propagação de fogo em grids 450×450 (30m/pixel)\n"
        "é computacionalmente caro. A CNN prediz a carga computacional por célula\n"
        "para balancear o processamento distribuído.\n\n"
        "Pipeline completo:\n"
        "  Dados LANDFIRE/NIFC  →  Preprocessing  →  elmfire (labels)\n"
        "  →  CNN (load prediction)  →  Simulador C sequencial (baseline)",
        Inches(0.5), Inches(1.3), Inches(8.5), Inches(5.5),
        fontsize=15, color=COR_TEXTO)

    # Diagrama simplificado
    boxes = [
        ("LANDFIRE\n+ NIFC", COR_TITULO),
        ("Semantic\n450×450×14", COR_DESTAQUE),
        ("elmfire\n(labels TOA)", PptxRGB(0x2E,0x7D,0x32)),
        ("CNN\n(load pred)", PptxRGB(0x6A,0x1B,0x9A)),
        ("Simulador C\n(baseline)", PptxRGB(0x01,0x57,0x9B)),
    ]
    for i, (txt, col) in enumerate(boxes):
        lx = Inches(9.3 + 0) + i * Inches(0)
        lx = Inches(9.0)
        ly = Inches(1.4 + i * 1.1)
        add_rect(s, lx, ly, Inches(3.8), Inches(0.85), col)
        add_title_box(s, txt, lx, ly + Inches(0.1), Inches(3.8), Inches(0.75),
                      fontsize=13, color=PptxRGB(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        if i < 4:
            add_body_text(s, "↓", lx + Inches(1.7), ly + Inches(0.82), Inches(0.5), Inches(0.3),
                          fontsize=14, color=COR_TITULO, bold=True)

    # ── SLIDE 4: Ambiente ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Configuração do Ambiente", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=26, color=PptxRGB(0xFF,0xFF,0xFF))

    env_items = [
        ("🖥 GPU", "NVIDIA RTX 4090 — 24 GB VRAM"),
        ("⚡ CUDA", "CUDA 12.1 + PyTorch 2.5.1 + cu121"),
        ("🔧 Compiler", "GCC 11 + OpenMP 5.0 + MPI (mpifort)"),
        ("🗺 GeoEspacial", "rasterio 1.5 + pyproj 3.7 + GDAL 3.8"),
        ("🔥 Simulador", "elmfire 2025.1002 (Fortran + MPI, compilado)"),
        ("📐 C Sequencial", "GCC -O2 -fopenmp (Rothermel + Dijkstra)"),
    ]
    for i, (icon_label, detail) in enumerate(env_items):
        col = i // 3
        row = i % 3
        lx = Inches(0.4 + col * 6.5)
        ly = Inches(1.4 + row * 1.85)
        add_rect(s, lx, ly, Inches(6.0), Inches(1.6), PptxRGB(0xFF,0xFF,0xFF), COR_TITULO)
        add_title_box(s, icon_label, lx + Inches(0.15), ly + Inches(0.1),
                      Inches(5.5), Inches(0.5), fontsize=14, color=COR_DESTAQUE)
        add_body_text(s, detail, lx + Inches(0.15), ly + Inches(0.6),
                      Inches(5.5), Inches(0.8), fontsize=13, color=COR_TEXTO)

    # ── SLIDE 5: Dataset ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Dataset — 68 Cenários de Incêndios 2016", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=26, color=PptxRGB(0xFF,0xFF,0xFF))

    add_body_text(s,
        "Fontes: NIFC Open Data (4.410 perímetros) + LANDFIRE 2016 (ArcGIS ImageServer)\n\n"
        "Critérios de seleção (per paper):\n"
        "  • Large-scale: ≥ 300 acres (USDA/NWCG definition)\n"
        "  • Apenas perímetro final (latest=Y) por incêndio\n"
        "  • Localização: EUA continentais (CONUS)\n\n"
        "Divisão:\n"
        "  • Treino: 51 cenários (aleatórios, classes E/F/G)\n"
        "  • Teste:  17 cenários (classe F: 1.000–4.999 acres)",
        Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5),
        fontsize=14, color=COR_TEXTO)

    # Stats boxes
    stats = [("68", "Cenários totais"), ("51", "Treino"), ("17", "Teste (Classe F)"),
             ("437", "Fires válidos no NIFC"), ("4.410", "Perímetros brutos")]
    for i, (num, label) in enumerate(stats):
        ly = Inches(1.4 + i * 1.1)
        add_rect(s, Inches(6.3), ly, Inches(2.2), Inches(0.9), COR_DESTAQUE)
        add_title_box(s, num, Inches(6.3), ly, Inches(2.2), Inches(0.55),
                      fontsize=24, color=PptxRGB(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        add_body_text(s, label, Inches(6.3), ly + Inches(0.5), Inches(2.2), Inches(0.45),
                      fontsize=11, color=PptxRGB(0xFF,0xFF,0xFF), bold=True)

    # ── SLIDE 6: Camadas Semânticas ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Pipeline de Preprocessing — 450×450×14", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=26, color=PptxRGB(0xFF,0xFF,0xFF))

    add_image_safe(s, VIZ / "semantic_S001.png",
                   Inches(0.3), Inches(1.2), Inches(9.0))
    add_body_text(s,
        "14 camadas por cenário\n\n"
        "8 Landscape (LANDFIRE):\n"
        "Elevação, Aspect, Slope,\n"
        "Fuel Model, Canopy Cover,\n"
        "Canopy Height, CBH, CBD\n\n"
        "6 Weather (sintético):\n"
        "Temperatura, Umidade,\n"
        "Cloud Cover, Precipitação,\n"
        "Wind Dir, Wind Speed\n\n"
        "Resolução: 30m/pixel\n"
        "Shape: (450, 450, 14)\n"
        "Normalizado: [0, 1]",
        Inches(9.5), Inches(1.2), Inches(3.5), Inches(6.0),
        fontsize=12, color=COR_TITULO)

    # ── SLIDE 7: elmfire Labels ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Geração de Labels — elmfire (Time of Arrival)", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=24, color=PptxRGB(0xFF,0xFF,0xFF))

    add_image_safe(s, VIZ / "elmfire_toa_S001.png",
                   Inches(0.3), Inches(1.2), Inches(7.5))
    add_body_text(s,
        "elmfire (v2025.1002)\n"
        "Algorithms:\n"
        "• Rothermel (1972) ROS\n"
        "• Huygens principle\n\n"
        "Input por cenário:\n"
        "• 8 GeoTIFFs landscape\n"
        "• 7 rasters de clima\n"
        "• Ponto de ignição (lat/lon)\n\n"
        "Output:\n"
        "• time_of_arrival.bil → .tif\n"
        "• Shape: (450, 450)\n"
        "• Valor: segundos desde ignição\n\n"
        "68 execuções em batch\n"
        "Paralelizado: 4 workers",
        Inches(8.0), Inches(1.3), Inches(5.0), Inches(5.8),
        fontsize=12, color=COR_TITULO)

    # ── SLIDE 8: CNN ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "CNN de Predição de Carga — LoadPredictorCNN", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=24, color=PptxRGB(0xFF,0xFF,0xFF))

    # Arquitetura boxes
    arch = [
        ("Input\n(14, W, W)", COR_TITULO),
        ("Conv2d(14→32)\nBN + ReLU", PptxRGB(0x01,0x57,0x9B)),
        ("Conv2d(32→64)\nBN + ReLU", PptxRGB(0x01,0x57,0x9B)),
        ("Conv2d(64→128)\nBN + ReLU", PptxRGB(0x01,0x57,0x9B)),
        ("Flatten\nFC(256)+ReLU", PptxRGB(0x2E,0x7D,0x32)),
        ("FC(1)\nOutput", COR_DESTAQUE),
    ]
    for i, (txt, col) in enumerate(arch):
        lx = Inches(0.3 + i * 2.18)
        add_rect(s, lx, Inches(1.4), Inches(2.0), Inches(1.1), col)
        add_title_box(s, txt, lx, Inches(1.45), Inches(2.0), Inches(1.0),
                      fontsize=12, color=PptxRGB(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        if i < 5:
            add_body_text(s, "→", lx + Inches(2.0), Inches(1.77), Inches(0.2), Inches(0.4),
                          fontsize=16, color=COR_TITULO, bold=True)

    # Sweep chart
    add_image_safe(s, VIZ / "sweep_mse.png",
                   Inches(0.3), Inches(2.7), Inches(8.0))

    add_body_text(s,
        "Sweep Atomic Sizes:\n"
        "9×9   → MSE 0.001498\n"
        "12×12 → MSE 0.001540\n"
        "15×15 → MSE 0.001488\n"
        "18×18 → MSE 0.001455\n"
        "21×21 → MSE 0.001448 ✓\n\n"
        "10.3M amostras treino\n"
        "3.4M amostras teste\n"
        "RTX 4090 — ~400s/época",
        Inches(8.5), Inches(2.7), Inches(4.5), Inches(4.5),
        fontsize=13, color=COR_TITULO)

    # ── SLIDE 9: Simulador C ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, COR_BG_SLIDE)
    add_rect(s, 0, 0, W, Inches(1.1), COR_TITULO)
    add_title_box(s, "Simulador Sequencial em C — Baseline", Inches(0.4), Inches(0.15), W, Inches(0.9),
                  fontsize=26, color=PptxRGB(0xFF,0xFF,0xFF))

    add_image_safe(s, VIZ / "simulator_S001_ws8_wd270.png",
                   Inches(0.3), Inches(1.2), Inches(8.5))

    add_body_text(s,
        "Algoritmo:\n"
        "• Rothermel (1972) ROS\n"
        "• Huygens (8-vizinhança)\n"
        "• Dijkstra (min-heap)\n\n"
        "Compilação:\n"
        "gcc -O2 -fopenmp\n\n"
        "Benchmark (68 cenários):\n"
        "• Médio:  26.6 ms\n"
        "• Mínimo: 25.2 ms\n"
        "• Máximo: 32.1 ms\n"
        "• Desvpad: 1.4 ms\n\n"
        "→ Este é o BASELINE\n"
        "   para as próximas\n"
        "   entregas (OpenMP/CUDA)",
        Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.8),
        fontsize=13, color=COR_TITULO)

    # ── SLIDE 10: Próximos Passos ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, PptxRGB(0x0D,0x1B,0x2A))
    add_rect(s, 0, 0, W, Inches(1.1), COR_DESTAQUE)
    add_title_box(s, "Próximos Passos — Entrega 2", Inches(0.4), Inches(0.1), W, Inches(0.9),
                  fontsize=28, color=PptxRGB(0xFF,0xFF,0xFF))

    steps = [
        ("🔀 OpenMP", "Paralelização do simulador C com #pragma omp\nSpeedup esperado: 8-16× (CPU multi-core)"),
        ("⚡ CUDA",   "Versão GPU do simulador\nSpeedup esperado: 50-200× vs. sequencial"),
        ("⚖️ Load Balancing", "Usar CNN para particionar a grade\nMinimizar waiting time entre nós"),
        ("📊 Benchmark Final", "Comparar: sequencial vs. OpenMP vs. CUDA\n38.5% de redução (paper) como meta"),
    ]
    for i, (title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        lx = Inches(0.5 + col * 6.4)
        ly = Inches(1.4 + row * 2.7)
        add_rect(s, lx, ly, Inches(6.0), Inches(2.4), PptxRGB(0x1A,0x3A,0x5C))
        add_title_box(s, title, lx + Inches(0.2), ly + Inches(0.1),
                      Inches(5.6), Inches(0.7), fontsize=16, color=COR_DESTAQUE)
        add_body_text(s, desc, lx + Inches(0.2), ly + Inches(0.8),
                      Inches(5.6), Inches(1.4), fontsize=13, color=PptxRGB(0xDD,0xDD,0xDD))

    out = OUT_DIR / "apresentacao_entrega1.pptx"
    prs.save(out)
    print(f"Apresentação salva: {out}")
    return out


if __name__ == "__main__":
    print("Gerando relatório DOCX...")
    build_report()
    print("\nGerando apresentação PPTX...")
    build_pptx()
    print("\nConcluído!")
